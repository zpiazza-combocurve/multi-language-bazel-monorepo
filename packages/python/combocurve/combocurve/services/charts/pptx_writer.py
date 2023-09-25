from typing import BinaryIO, Optional, Iterable
from io import BytesIO

from pptx import Presentation
from pptx.util import Mm, Inches
from pptx.slide import Slide
from pptx.shapes.picture import Picture

from .document_writer import ChartsDocumentWriter

BLANK_LAYOUT_INDEX = 6


class ChartsPptxWriter(ChartsDocumentWriter):
    margin = Inches(0.25)

    def __init__(self, file: BinaryIO, rows: int, columns: int, is_landscape: bool = True):
        super().__init__(file, rows, columns, is_landscape)

        if self.is_landscape:
            self.page_width, self.page_height = Mm(297), Mm(210)
        else:
            self.page_width, self.page_height = Mm(210), Mm(297)
        self.inner_width = self.page_width - 2 * self.margin
        self.inner_height = self.page_height - 2 * self.margin

        self.image_width = self.inner_width / columns
        self.image_height = self.inner_height / rows
        self.presentation = Presentation()
        self.presentation.slide_width = self.page_width
        self.presentation.slide_height = self.page_height
        self.blank_slide_layout = self.presentation.slide_layouts[BLANK_LAYOUT_INDEX]
        self.current_slide: Optional[Slide] = None

    def draw_image(self, image_file: BinaryIO, row: int, column: int):
        if self.current_slide is None:
            raise AssertionError(
                'Attempted to draw image before getting first page. Remember to call `next_page` first.')

        left = self.margin + column * self.image_width
        top = self.margin + row * self.image_height
        self.current_slide.shapes.add_picture(image_file, left, top, self.image_width, self.image_height)

    def next_page(self):
        self.current_slide = self.presentation.slides.add_slide(self.blank_slide_layout)

    def close(self):
        self.presentation.save(self.file)

    @staticmethod
    def merge_docs(document_files: Iterable[BinaryIO], output_file: BinaryIO):
        outputPresentation = Presentation()
        blank_layout = outputPresentation.slide_layouts[BLANK_LAYOUT_INDEX]

        for doc in document_files:
            inputPresentation = Presentation(doc)
            outputPresentation.slide_width = inputPresentation.slide_width
            outputPresentation.slide_height = inputPresentation.slide_height
            for inputSlide in inputPresentation.slides:
                outputSlide = outputPresentation.slides.add_slide(blank_layout)
                for shape in inputSlide.shapes:
                    if not isinstance(shape, Picture):
                        continue
                    with BytesIO(shape.image.blob) as image_file:
                        outputSlide.shapes.add_picture(image_file, shape.left, shape.top, shape.width, shape.height)

        outputPresentation.save(output_file)
